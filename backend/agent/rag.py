"""
RAG (Retrieval-Augmented Generation) Module for kGPT.

Handles document ingestion, vector storage with Chroma, embedding generation,
and RetrievalQA chain construction.
"""

import os
from typing import List

from dotenv import load_dotenv
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    CSVLoader,
    TextLoader,
    WebBaseLoader,
)
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain.chains import RetrievalQA
from langchain_core.documents import Document

load_dotenv()

CHROMA_PERSIST_DIR: str = os.getenv("CHROMA_PERSIST_DIR", "./vectorstore")

_embeddings_instance: HuggingFaceEmbeddings = None
_vectorstore_instance: Chroma = None

_text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
)

class _ExcelLoader:
    """Minimal .xlsx/.xlsm loader (uses openpyxl) -> one text Document."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> List[Document]:
        import openpyxl

        wb = openpyxl.load_workbook(self.file_path, data_only=True, read_only=True)
        blocks = []
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                blocks.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
        text = "\n\n".join(blocks)
        return [Document(page_content=text, metadata={"source": self.file_path})]


class _PptxLoader:
    """Minimal .pptx loader (uses python-pptx) -> one text Document."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> List[Document]:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            parts.append(line)
            if parts:
                slides.append(f"Slide {i}:\n" + "\n".join(parts))
        text = "\n\n".join(slides)
        return [Document(page_content=text, metadata={"source": self.file_path})]


_LOADER_MAP = {
    "pdf": PyPDFLoader,
    "docx": Docx2txtLoader,
    "csv": CSVLoader,
    "txt": TextLoader,
    "text": TextLoader,
    "md": TextLoader,
    "xlsx": _ExcelLoader,
    "pptx": _PptxLoader,
}


def get_embeddings() -> HuggingFaceEmbeddings:
    """
    Returns a singleton HuggingFaceEmbeddings instance using the
    all-MiniLM-L6-v2 model.
    """
    global _embeddings_instance
    if _embeddings_instance is None:
        _embeddings_instance = HuggingFaceEmbeddings(
            model_name="all-MiniLM-L6-v2",
            model_kwargs={"device": "cpu"},
            encode_kwargs={"normalize_embeddings": True},
        )
    return _embeddings_instance


def get_vectorstore() -> Chroma:
    """
    Returns a singleton Chroma vectorstore instance. Creates the persist
    directory and collection if they don't exist yet.
    """
    global _vectorstore_instance
    if _vectorstore_instance is None:
        os.makedirs(CHROMA_PERSIST_DIR, exist_ok=True)
        _vectorstore_instance = Chroma(
            collection_name="kgpt_documents",
            embedding_function=get_embeddings(),
            persist_directory=CHROMA_PERSIST_DIR,
        )
    return _vectorstore_instance


def ingest_document(file_path: str, file_type: str) -> int:
    """
    Loads a document from disk, splits it into chunks, and adds it to the
    vectorstore.

    Args:
        file_path: Absolute or relative path to the document file.
        file_type: One of 'pdf', 'docx', 'csv', 'txt', 'text', 'md'.

    Returns:
        The number of chunks added to the vectorstore.

    Raises:
        ValueError: If the file type is unsupported.
        FileNotFoundError: If the file does not exist.
    """
    file_type = file_type.strip().lower()

    if file_type not in _LOADER_MAP:
        raise ValueError(
            f"Unsupported file type: '{file_type}'. "
            f"Supported types: {list(_LOADER_MAP.keys())}"
        )

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    loader_cls = _LOADER_MAP[file_type]
    loader = loader_cls(file_path)
    documents: List[Document] = loader.load()

    chunks = _text_splitter.split_documents(documents)

    vectorstore = get_vectorstore()
    vectorstore.add_documents(chunks)

    return len(chunks)


def ingest_url(url: str) -> int:
    """
    Loads content from a web URL, splits it into chunks, and adds it to the
    vectorstore.

    Args:
        url: The web URL to scrape and ingest.

    Returns:
        The number of chunks added to the vectorstore.
    """
    loader = WebBaseLoader(url)
    documents: List[Document] = loader.load()

    chunks = _text_splitter.split_documents(documents)

    vectorstore = get_vectorstore()
    vectorstore.add_documents(chunks)

    return len(chunks)


def query_documents(question: str, k: int = 3) -> List[Document]:
    """
    Performs a similarity search against the vectorstore.

    Args:
        question: The natural-language query.
        k: Number of most-similar documents to return.

    Returns:
        A list of LangChain Document objects most relevant to the query.
    """
    vectorstore = get_vectorstore()
    results = vectorstore.similarity_search(question, k=k)
    return results


def get_rag_chain(llm):
    """
    Constructs a RetrievalQA chain that combines the vectorstore retriever
    with the provided LLM.

    Args:
        llm: A LangChain chat model instance.

    Returns:
        A RetrievalQA chain configured with the 'stuff' chain type.
    """
    vectorstore = get_vectorstore()
    retriever = vectorstore.as_retriever(
        search_type="similarity",
        search_kwargs={"k": 3},
    )

    rag_chain = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=retriever,
        return_source_documents=True,
        verbose=False,
    )

    return rag_chain